import glob
import pandas as pd

path_list = glob.glob('application/dash/assets/PowerpointTemplates/*.pptx')
pptx_templates = [file[44:] for file in path_list]
print(pptx_templates)

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.chart.chart import Chart



# Run a loop here that goes through each slide and prints the objects / other stuff
prs = Presentation('application/dash/assets/PowerpointTemplates/ModernAngles.pptx')

for slide in prs.slides:
    print("Slide: ",slide.name, slide.slide_id)
    for shape in slide.shapes:
        print(shape.name, shape.shape_id, shape.has_chart)

# Function to put a data-frame in a chart 

def find_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape

# I want to make this better by clarifying what classes each of these objects need to be, and maybe some extra code to confirm that rows/columns of data-frame are strings
def dataframe_to_chart(chart, df, format = '0%;0%'):
    chart_data = CategoryChartData(number_format = format)
    chart_data.categories = df.index
    for srs in df.columns:
        chart_data.add_series(srs, df[srs])
    chart.replace_data(chart_data)


test_df = pd.DataFrame([[0.25,0.35,0.45,0.3],[0.35,0.25,0.5,0.6]],index=['Alpha','Beta'],columns = ['t1','t2','t3','t4'])

dataframe_to_chart(find_shape(prs.slides[1],6).chart,test_df)
prs.save('application/dash/assets/PowerpointTemplates/test.pptx')