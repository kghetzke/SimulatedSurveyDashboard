
import dash_html_components as html
from dash.dependencies import Input, Output, State, ALL
from application.dash.app import app
from datetime import datetime
import pandas as pd
import plotly.express as px
import dash_core_components as dcc
import application.app_settings as config
from application.dash.components.datasteps import get_data, get_meta

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.chart.chart import Chart
import io

"""Some helper functions for using powerpoint -> 
    I wrote these to make better use of the Python-pptx package
"""
def find_shape(slide, shape_id):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape

def dataframe_to_chart(chart, df, format = '0%;0%'):
    """ This replaces data in an exisiting powerpoint chart object with a pandas dataframe
        It is necessary to have the row-indices and column names be strings; 
    """
    chart_data = CategoryChartData(number_format = format)
    chart_data.categories = df.index
    for srs in df.columns:
        chart_data.add_series(srs, df[srs])
    chart.replace_data(chart_data)




#New Callbacks:
#First, load the data in memory using my fetch_data functions in the datasteps.py script
df = get_data()
vars, vals = get_meta()

#Let's do a callback that creates the dynamic list of dropdowns from the series json stored in memory
@app.callback(
    [Output(component_id="dropdown_container", component_property = "children")],
    [Input(component_id="queried_series", component_property = "data")]
    )
def render_dropdowns(series_logic):
    df = pd.read_json(series_logic)
    dummy_count = len(df)
    dynamic_dropdowns =     [html.Div(children = [
        html.Div('Series '+str(k+1), style = {'width': "25%"}),
        dcc.Dropdown(id = {'type': 'alpha', 'index': k}, options=[{'label': vars[i], 'value': i} for i in list(vars.keys())[3:]], 
                     multi = False, value=df['alpha'][k], style={'width': "95%"}),
        dcc.Dropdown(id = {'type': 'beta', 'index': k}, options=[{'label': vals['company'][i], 'value': i} for i in vals['company']], 
                     multi = False, value=df['beta'][k], style={'width': "95%"}),
        dcc.Dropdown(id = {'type': 'gamma', 'index': k}, options=[{'label': vals['audience'][i], 'value': i} for i in vals['audience']], 
                     multi = False, value=df['gamma'][k], style={'width': "95%"})], style={'display': 'flex'}) for k in range(dummy_count)]
    return [dynamic_dropdowns]

#This callback does two different things depending on what is input:
#When "add series" is clicked, this appends a new series to the series_logic json
#When a dropdown changes, this updates the series_logic json
@app.callback(
    [Output('queried_series','data')],[Output('add_button','n_clicks')],
    [Input(component_id="add_button", component_property = "n_clicks")],
    [Input({'type': 'alpha', 'index': ALL},'value')],
    [Input({'type': 'beta', 'index': ALL},'value')],
    [Input({'type': 'gamma', 'index': ALL},'value')],
    [State('queried_series','data')],
    prevent_initial_call = True
    )
def update_logic(n,alphas, betas, gammas, df_json):
    df_logic = pd.read_json(df_json)
    if n==1:
        new_row = pd.DataFrame([[list(vars.keys())[5],1,1]], columns = ['alpha','beta','gamma'])
        df_logic = df_logic.append(new_row).reset_index(drop=True)
    else:
        for idx, val in enumerate(alphas):
            df_logic['alpha'][idx] = val
        for idx, val in enumerate(betas):
            df_logic['beta'][idx] = val
        for idx, val in enumerate(gammas):
            df_logic['gamma'][idx] = val
    return df_logic.to_json(), 0

##Update a graph based on the json (and store the df in memory)
@app.callback(
    [Output(component_id= 'main_graph', component_property = 'figure'),Output(component_id= 'df_memory', component_property = 'data')],
    [Input(component_id='queried_series',component_property='data')],
    )
def update_graph(logic_json):

    #read the logic stored in json
    logic_df = pd.read_json(logic_json)
    dummy_count = len(logic_df)

    # do some data-steps with a copy of our data-frame
    dff = df.copy()
    graph_df = dff[['month','Month','company','Company','audience','Audience',logic_df['alpha'][0]]][(dff['audience']==logic_df['gamma'][0]) & (dff['company']==logic_df['beta'][0])]
    graph_df = graph_df.rename(columns = {logic_df['alpha'][0]: 'Value'})
    #graph_df = graph_df.dropna(subset = ['Value'], inplace=True)
    graph_df['Series'] = 'Series 1'
    graph_df['Series_Idx'] = 1
    graph_df['Metric'] = logic_df['alpha'][0]

    if dummy_count>1:
        for i in range(1,dummy_count):
            new_rows = dff[['month','Month','company','Company','audience','Audience',logic_df['alpha'][i]]][(dff['audience']==logic_df['gamma'][i]) & (dff['company']==logic_df['beta'][i])]
            new_rows = new_rows.rename(columns = {logic_df['alpha'][i]: 'Value'})
            #new_rows = new_rows.dropna(subset = ['Value'], inplace=True)
            new_rows['Series'] = 'Series ' + str(i+1)
            new_rows['Series_Idx'] = i+1
            new_rows['Metric'] = logic_df['alpha'][i]
            graph_df = graph_df.append(new_rows).reset_index(drop = True)

    # render the figure using plotly-express library
    fig = px.line(graph_df, x="Month", y='Value', color='Series', labels = {'Value': 'Value (% as decimal)'})
    
    # return the figure (which will be sent to the 'my_graph' placeholder via the callback) 
    # and the graph_df as json (which will be sent to the "df_memory" dcc.store() object)
    return [fig, graph_df[['Series','month','Month','Company','Audience','Metric','Value']].to_json()]


#Download the df stored in memory when clicking the download button
#I need to read the df_memory in State(), not Input() otherwise I'll trigger a download every time the user updates the dropdowns
@app.callback(
    [Output(component_id= "download_csv", component_property = "data")],
    [Input("download_data_button", "n_clicks",)],
    [State("df_memory","data")],
    prevent_initial_call = True,
    )
def download_data(n_clicks, df_json):
    dndf = pd.read_json(df_json)
    return [dcc.send_data_frame(dndf.to_csv, "data.csv", index=False)]

#Take the DF Stored in Memory and Download a Powerpoint.pptx file
#I need to read the df_memory in State(), not Input() otherwise I'll trigger a download every time the user updates the dropdowns
@app.callback(
    [Output(component_id= "download_pptx", component_property = "data")],
    [Input("download_pptx_button", "value",)],
    [State("df_memory","data")],
    prevent_initial_call = True,
    )
def download_pptx(file, df_json):
    """This uses the py-pptx package to create a presentation object from an existing presentation, swap data in chart, the save in an ioBytes object to download"""
    dndf = pd.read_json(df_json)
    ndf = dndf[['Series','month','Value']]
    ndf = ndf.pivot(index='month',columns = 'Series',values='Value')
    ndf.index = ['January','February','March','April','May','June','July','August','September','October','November','December']
    prs = Presentation('application/dash/assets/PowerpointTemplates/' + file)
    dataframe_to_chart(find_shape(prs.slides[1],6).chart,ndf)
    file_obj = io.BytesIO()
    prs.save(file_obj)
    file_obj.seek(0)
    return [dcc.send_bytes(file_obj.read(), "DataChart.pptx")]